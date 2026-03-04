"""Engine that converts PPTX slides to images and uses VLM for text extraction."""

from __future__ import annotations

import logging
import subprocess
import tempfile
from pathlib import Path
from typing import Any, BinaryIO, Callable, Dict, List, Optional

from ..base import ParsedBlock, ParsedDocument, ParsedPage, PdfParseOptions

logger = logging.getLogger(__name__)

DEFAULT_PROMPT = (
    "Read this page from a semiconductor equipment technical document "
    "and extract all content as Markdown.\n\n"
    "Rules:\n"
    "- Output ONLY valid Markdown (no LaTeX, no $...$ or $...$)\n"
    "- Convert mathematical formulas to plain text or Unicode symbols\n"
    "- Preserve tables using Markdown table syntax (| column |)\n"
    "- Keep headings with # syntax\n"
    "- Preserve lists and bullet points\n"
    "- Preserve numbered steps and procedure sequences exactly (e.g. Step 1, 1), (1))\n"
    "- Keep all equipment/part names verbatim "
    "(e.g. APC valve, MFC, Chuck Motor, EFEM Robot, Load Lock)\n"
    "- Extract text visible in photos, diagrams, and labels "
    "(arrows, callouts, part labels) as [Image text: ...]\n"
    "- Preserve WARNING, CAUTION, and NOTE blocks with their markers\n"
    "- Keep Korean-English mixed text as-is without translating\n"
    "- Do not summarize or omit any text\n"
    "- Do not add explanations or comments\n"
    "- NEVER output repeated characters like ||||, &&&&, ----, ====, etc."
)


class PptxVlmEngine:
    """PPTX -> slide images -> VLM text extraction engine.

    Pipeline:
        1. Write PPTX bytes to temp file
        2. Convert PPTX -> PDF via libreoffice --headless
        3. Render PDF pages -> PIL images via pdf2image
        4. Each image -> VLM generate() -> text
        5. Return ParsedDocument (slide = page)

    Fallback: If libreoffice conversion fails, extract text directly
    from PPTX using python-pptx (text-only, no layout/images).
    """

    content_type: str = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

    def __init__(
        self,
        *,
        vlm_client: Any | None = None,
        vlm_factory: Callable[[PdfParseOptions], Any] | None = None,
        default_prompt: str = DEFAULT_PROMPT,
    ) -> None:
        self.vlm_client = vlm_client
        self.vlm_factory = vlm_factory
        self.default_prompt = default_prompt

    def _pptx_to_images(self, pptx_bytes: bytes) -> List[Any]:
        """PPTX -> slide images via libreoffice + pdf2image."""
        try:
            from pdf2image import convert_from_path
        except ImportError as exc:
            raise ImportError(
                "pdf2image is required for PPTX VLM parsing. "
                "Install with: pip install pdf2image"
            ) from exc

        with tempfile.TemporaryDirectory() as tmpdir:
            tmpdir_path = Path(tmpdir)
            pptx_path = tmpdir_path / "input.pptx"
            pptx_path.write_bytes(pptx_bytes)

            try:
                subprocess.run(
                    [
                        "libreoffice",
                        "--headless",
                        "--convert-to", "pdf",
                        "--outdir", str(tmpdir_path),
                        str(pptx_path),
                    ],
                    capture_output=True,
                    text=True,
                    check=True,
                    timeout=120,
                )
            except (subprocess.CalledProcessError, FileNotFoundError, subprocess.TimeoutExpired) as exc:
                logger.warning(
                    "libreoffice PPTX->PDF failed: %s. Falling back to text extraction.", exc,
                )
                return []

            pdf_path = tmpdir_path / "input.pdf"
            if not pdf_path.exists():
                logger.warning("libreoffice did not produce PDF, falling back")
                return []

            images = convert_from_path(str(pdf_path), dpi=200)
            return images

    def _extract_text_fallback(self, pptx_bytes: bytes) -> List[ParsedPage]:
        """python-pptx text-only extraction (fallback)."""
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise ImportError(
                "python-pptx is required for PPTX text fallback. "
                "Install with: pip install python-pptx"
            ) from exc

        import io
        prs = Presentation(io.BytesIO(pptx_bytes))
        pages: List[ParsedPage] = []

        for slide_idx, slide in enumerate(prs.slides, start=1):
            texts: List[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            texts.append(text)
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells)
                        if row_text.strip(" |"):
                            texts.append(row_text)

            page_text = "\n".join(texts)
            pages.append(ParsedPage(number=slide_idx, text=page_text))

        return pages

    def run(self, file: BinaryIO, options: Optional[PdfParseOptions] = None) -> ParsedDocument:
        """Parse PPTX file via VLM.

        Args:
            file: PPTX binary stream
            options: Parse options

        Returns:
            ParsedDocument (slide = page)
        """
        opts = options or PdfParseOptions()

        if self.vlm_client is None and self.vlm_factory is not None:
            self.vlm_client = self.vlm_factory(opts)

        if hasattr(file, "seek"):
            try:
                file.seek(0)
            except Exception:
                pass
        pptx_bytes = file.read()

        images = self._pptx_to_images(pptx_bytes)

        if not images:
            logger.info("Using python-pptx text fallback")
            pages = self._extract_text_fallback(pptx_bytes)
            blocks = [ParsedBlock(text=p.text, page=p.number, label="slide") for p in pages]
            return ParsedDocument(
                pages=pages, blocks=blocks, tables=[], figures=[],
                metadata={"parser": "pptx_vlm", "method": "text_fallback"},
                errors=["libreoffice conversion failed, used text-only fallback"],
                content_type=self.content_type,
            )

        if self.vlm_client is None:
            raise ImportError("vlm_client is required for PptxVlmEngine VLM parsing")

        pages: List[ParsedPage] = []
        blocks: List[ParsedBlock] = []

        for index, image in enumerate(images, start=1):
            prompt = opts.vlm_prompt or self.default_prompt
            infer_kwargs: Dict[str, Any] = {}
            if opts.vlm_max_new_tokens is not None:
                infer_kwargs["max_tokens"] = opts.vlm_max_new_tokens
            if opts.vlm_temperature is not None:
                infer_kwargs["temperature"] = opts.vlm_temperature

            text = self.vlm_client.generate(image=image, prompt=prompt, **infer_kwargs)
            pages.append(ParsedPage(number=index, text=text))
            blocks.append(ParsedBlock(text=text, page=index, label="slide"))

        metadata = {
            "parser": "pptx_vlm",
            "method": "libreoffice_vlm",
            "vlm_model": opts.vlm_model,
            "vlm_prompt": opts.vlm_prompt or self.default_prompt,
            "total_slides": len(pages),
        }

        return ParsedDocument(
            pages=pages, blocks=blocks, tables=[], figures=[],
            metadata=metadata, errors=[], content_type=self.content_type,
        )


__all__ = ["PptxVlmEngine"]
